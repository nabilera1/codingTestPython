from pptx import Presentation


# 새 프레젠테이션 생성
prs = Presentation()

# 제목 슬라이드 추가
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "2025 체육대회 행사 기획안"
subtitle.text = ""

# 슬라이드 내용 리스트
slides_content = [
    ("🏟️ 체육대회 개최 안내", ["📅 개최일: 2025년 5월 15일 (금요일, 날씨 맑음 기원! ☀️)",
                           "📍 장소: 학교 운동장 및 체육관",
                           "🎉 대상: 전 학년 학생 및 교직원 (함께 즐기는 축제의 장!)"]),

    ("💃 개막식 & 퍼레이드", ["🎶 신나는 음악과 함께 체육대회 개막!",
                          "🏳️‍🌈 각 반별 입장 퍼레이드 (컨셉은 자유! 😆)",
                          "🎤 교장 선생님의 한마디 & 축하 공연"]),

    ("🔥 체육대회 주요 경기", ["🏃‍♂️ 100m 달리기 – 스피드 최강자는 누구?!",
                          "🏋️‍♀️ 줄다리기 – 반 대항 팀워크 대결!",
                          "⚽ 축구 & 피구 – 승리를 위한 불꽃 경쟁!",
                          "🏀 농구 & 발야구 – 최고의 스포츠 스타 탄생?",
                          "🎭 번외 경기 – 교사 vs 학생 이색 경기!"]),

    ("🍱 점심 & 간식 타임", ["🍙 도시락 & 간식 제공 (에너지 충전 완료!)",
                           "🍧 시원한 음료와 아이스크림 (더위를 날려버려~)"]),

    ("🏆 체육대회 이벤트", ["🎤 응원전 – 반별 창의적인 응원 퍼포먼스!",
                         "🎯 행운의 럭키드로우 – 깜짝 경품 증정 🎁",
                         "📸 포토존 운영 – 특별한 순간을 사진으로 남기자!"]),

    ("🎖️ 시상식 & 폐막식", ["🏅 종합 우승 반 시상!",
                           "🥇 MVP 선수 선정 – 체육대회의 히어로는 누구?",
                           "🎉 폐막 선언 & 단체 사진 촬영 📷"]),

    ("💬 추가 논의 및 의견 수렴", ["📢 더 재미있는 경기 아이디어 대환영!",
                               "💡 체육대회 운영 관련 의견 자유롭게 제안 가능!"])
]


# 각 슬라이드 추가
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # 제목 + 내용 레이아웃
    slide = prs.slides.add_slide(slide_layout)
    slide_title = slide.shapes.title
    slide_content = slide.placeholders[1]

    slide_title.text = title
    slide_content.text = "\n".join(content)

# 파일 저장
pptx_filename = "./2025 체육대회 행사 기획안1.pptx"
prs.save(pptx_filename)

# 파일 경로 반환
pptx_filename
